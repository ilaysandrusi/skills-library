#!/usr/bin/env python3
"""Four decks. The last two are the ones that keep the rule from collapsing in either direction.

  mac_fonts.pptx   the real failure: a body font bundled with macOS, plus a stray monospace. This
                   is what a deck built on a laptop and shipped to a Windows house PC looks like.
  portable.pptx    the same slides in fonts that exist on both platforms. Must be silent, or the
                   check is telling everyone their deck is broken.
  embedded.pptx    names the macOS font AND embeds it, which is one of the two correct answers to
                   this problem. A check that fires here punishes the fix, and a check that
                   punishes the fix gets switched off.
  korean_text.pptx sets no run font at all, so its Korean body text falls to the stock template's
                   theme — which declares a Windows-only Korean face. This one MUST fire. Without
                   it, "an inherited default only counts when the script is present" could be
                   implemented as "an inherited default never counts" and every test would still
                   be green.

Written wherever the caller says (a temp dir, in practice). Nothing built lands in the repo tree.
Needs python-pptx (CI installs it).
"""

from __future__ import annotations

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

BLANK = 6
BODY = ("412 consecutive patients; recurrence 12% versus 26% over a median 3.2 years, "
        "unchanged after adjustment for tumour size.")


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def box(slide, text, x, y, w, h, pt, font=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(pt)
    if font:
        run.font.name = font
    return tb


def build(path: Path, body_font: str, mono_font: str) -> None:
    prs = deck()
    for i in range(4):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, f"Adjunctive ablation halved local recurrence ({i + 1})",
            0.8, 0.8, 11.5, 1.1, 32, body_font)
        box(s, BODY, 0.8, 2.3, 11.5, 2.0, 20, body_font)
        box(s, "analysis/fit_model.R", 0.8, 5.0, 6.0, 0.5, 14, mono_font)
    prs.save(str(path))


# The presentation part is where a deck declares the fonts it carries. python-pptx has no API for
# it, so the fixture writes the element itself — which is also the honest way to test the exemption,
# because the detector's claim is about the XML and not about how the XML got there.
EMBEDDED = ('<p:embeddedFontLst><p:embeddedFont>'
            '<p:font typeface="Apple SD Gothic Neo" pitchFamily="34" charset="-127"/>'
            '</p:embeddedFont><p:embeddedFont>'
            '<p:font typeface="Menlo" pitchFamily="49" charset="0"/>'
            '</p:embeddedFont></p:embeddedFontLst>')


def add_embedded_font_list(path: Path) -> None:
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            data = zin.read(item)
            if item == "ppt/presentation.xml":
                xml = data.decode("utf-8")
                # CT_Presentation orders embeddedFontLst after notesSz.
                xml, n = re.subn(r"(<p:notesSz[^>]*/>)", r"\1" + EMBEDDED, xml, count=1)
                if n != 1:
                    raise SystemExit("fixture: could not find <p:notesSz/> to anchor the font list")
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def build_korean(path: Path) -> None:
    """No run font anywhere. The Korean text falls through to the theme's East-Asian fallback."""
    prs = deck()
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[BLANK])
        box(s, f"국소 재발률이 절반으로 줄었습니다 ({i + 1})", 0.8, 0.8, 11.5, 1.1, 32)
        box(s, "연속 환자 412명, 중앙 추적 3.2년에서 재발률은 12% 대 26%였습니다.",
            0.8, 2.3, 11.5, 2.0, 20)
    prs.save(str(path))


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    out.mkdir(parents=True, exist_ok=True)
    build(out / "mac_fonts.pptx", "Apple SD Gothic Neo", "Menlo")
    build(out / "portable.pptx", "Inter", "Noto Sans Mono")
    build(out / "embedded.pptx", "Apple SD Gothic Neo", "Menlo")
    add_embedded_font_list(out / "embedded.pptx")
    build_korean(out / "korean_text.pptx")
    print(f"wrote 4 decks into {out}")

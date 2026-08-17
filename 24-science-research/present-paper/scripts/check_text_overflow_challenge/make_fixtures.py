#!/usr/bin/env python3
"""One deck, and three recorded measurements of it.

The detector's whole argument is that the render already knows the answer, so the fixture supplies
a render — as `pdftotext -bbox` output, which is exactly what the detector reads. Recording the
measurement rather than producing a PDF keeps the card deterministic and free of poppler and
LibreOffice, and it exercises the same parser and the same comparison on the same geometry.

  deck.pptx        two slides. Slide 1 carries a filled block with a known rectangle.
  clean.xml        every line inside its block and clear of the foot of the slide.
  overflow.xml     one line whose bottom passes the block's bottom, and one that ends inside the
                   reserved band at the foot of the slide.
  short.xml        one page for a two-slide deck — a handout export, where page N is not slide N.
                   The detector must refuse to compare rather than compare the wrong things.

The coordinates in the XML are computed from the same constants that place the shapes, so the
fixture cannot drift away from the deck it describes.

Needs python-pptx (CI installs it).
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

PT = 72.0
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

# The block on slide 1, in inches. Everything the XML says is derived from these four numbers.
CARD_X, CARD_Y, CARD_W, CARD_H = 0.8, 1.5, 11.5, 2.0
CARD_TOP_PT = CARD_Y * PT
CARD_BOTTOM_PT = (CARD_Y + CARD_H) * PT
CARD_LEFT_PT = CARD_X * PT
CARD_RIGHT_PT = (CARD_X + CARD_W) * PT


def build_deck(path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(CARD_X), Inches(CARD_Y),
                               Inches(CARD_W), Inches(CARD_H))
    card.fill.solid()
    tf = card.text_frame
    tf.word_wrap = True
    tf.text = "Recurrence was 12% versus 26% over a median of 3.2 years."
    tf.paragraphs[0].runs[0].font.size = Pt(20)

    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(2.0))
    tb.text_frame.word_wrap = True
    tb.text_frame.text = "The tract was the route."
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    prs.save(str(path))


def line(x0: float, y0: float, x1: float, y1: float, text: str) -> str:
    return (f'<line xMin="{x0:.2f}" yMin="{y0:.2f}" xMax="{x1:.2f}" yMax="{y1:.2f}">'
            f"<word>{text}</word></line>")


def page(lines: list) -> str:
    return (f'<page width="{SLIDE_W_IN * PT:.2f}" height="{SLIDE_H_IN * PT:.2f}">'
            + "".join(lines) + "</page>")


HEAD = ('<?xml version="1.0" encoding="UTF-8"?>\n<html><body>\n')
TAIL = "\n</body></html>\n"


def write(path: Path, pages: list) -> None:
    path.write_text(HEAD + "\n".join(pages) + TAIL, encoding="utf-8")


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    out.mkdir(parents=True, exist_ok=True)
    build_deck(out / "deck.pptx")

    inside = line(CARD_LEFT_PT + 12, CARD_TOP_PT + 20, CARD_RIGHT_PT - 12, CARD_TOP_PT + 44,
                  "Recurrence was 12% versus 26%.")
    high_up = line(60, 100, 600, 130, "The tract was the route.")

    write(out / "clean.xml", [page([inside]), page([high_up])])

    # Bottom 0.10 in is reserved by default, so a line ending below (7.5 - 0.10) in is off-slide.
    past_block = line(CARD_LEFT_PT + 12, CARD_BOTTOM_PT - 12, CARD_RIGHT_PT - 12,
                      CARD_BOTTOM_PT + 48, "and 24.6% carried no reference standard at all")
    in_footer = line(60, (SLIDE_H_IN - 0.30) * PT, 600, (SLIDE_H_IN - 0.02) * PT,
                     "Supported by an institutional grant.")
    write(out / "overflow.xml", [page([inside, past_block]), page([high_up, in_footer])])

    write(out / "short.xml", [page([inside])])

    print(f"wrote deck.pptx + 3 recorded measurements into {out}")
